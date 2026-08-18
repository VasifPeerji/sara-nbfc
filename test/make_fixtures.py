#!/usr/bin/env python3
"""Generate real documents for the file-parser tests.

test_files.js builds its own fixtures in JavaScript so it runs anywhere with
no dependencies. This script adds the stronger check: files written by the
libraries real people use, so the parser is verified against genuine producer
output rather than against our own idea of what producers emit.

    python test/make_fixtures.py

Anything that cannot be produced on this machine is skipped with a note. The
Node suite treats every fixture as optional and says which ones it found.
"""

import shutil
import sys
import zipfile
from pathlib import Path

OUT = Path(__file__).resolve().parent / "fixtures"

# A sentence with characters that break naive parsers: an ampersand and a
# less-than that must survive XML escaping, and a status word that the
# renderer highlights, so the two features are exercised together.
SENTENCE = "Cross-check R&D spend < 5% against the compliance register."
MARKER = "ZEBRAFISH"          # unique token, proves text came from this file
TABLE = [["Site", "Status", "Reviewed"],
         ["Northam", "compliant", "2026-03-11"],
         ["Geraldton", "at risk", "2026-05-02"]]


def docx(path):
    """Hand-built WordprocessingML in a real ZIP.

    python-docx is not installed here, so the package is written directly.
    The XML is what Word emits for these constructs: paragraphs, a tab, a
    line break, escaped entities and a table.
    """
    def p(text):
        return f'<w:p><w:r><w:t xml:space="preserve">{text}</w:t></w:r></w:p>'

    rows = "".join(
        "<w:tr>" + "".join(
            f'<w:tc><w:p><w:r><w:t>{c}</w:t></w:r></w:p></w:tc>' for c in row
        ) + "</w:tr>"
        for row in TABLE
    )

    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        + p(f"{MARKER} quarterly review")
        + p("Cross-check R&amp;D spend &lt; 5% against the compliance register.")
        + '<w:p><w:r><w:t>Before</w:t><w:tab/><w:t>after a tab</w:t>'
          "<w:br/><w:t>after a break</w:t></w:r></w:p>"
        + f"<w:tbl>{rows}</w:tbl>"
        + p("Closing paragraph.")
        + "</w:body></w:document>"
    )

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Target="word/document.xml" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"/>'
        "</Relationships>"
    )

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
    return "hand-written OOXML in a real deflate ZIP"


def xlsx(path):
    import datetime
    import openpyxl
    wb = openpyxl.Workbook()

    first = wb.active
    first.title = "Sites"
    for row in TABLE:
        first.append(row)
    first["E1"] = MARKER                 # a gap before it, to test column mapping
    first["A5"] = SENTENCE
    first["B6"] = 1234.5
    first["C6"] = "=B6*2"               # a formula: the cached value must not
                                        # be confused with the formula text

    # Real typed cells. Excel stores these as bare numbers plus a format, so a
    # parser that reads only the value reports 46092 for a date and 0.075 for
    # a percentage. This sheet is here to make sure that never comes back.
    typed = wb.create_sheet("Typed")
    typed["A1"] = "Due"
    typed["B1"] = datetime.date(2026, 3, 11)
    typed["A2"] = "Meeting"
    typed["B2"] = datetime.datetime(2026, 12, 1, 14, 30)
    typed["A3"] = "Rate"
    typed["B3"] = 0.075
    typed["B3"].number_format = "0.0%"
    typed["A4"] = "Custom"
    typed["B4"] = datetime.date(2026, 5, 2)
    typed["B4"].number_format = "dd/mm/yyyy"

    second = wb.create_sheet("Notes")
    second["A1"] = "Second sheet reached"

    wb.save(path)
    return f"openpyxl {openpyxl.__version__}"


def pdf_layout(path):
    """Two columns and a table.

    A PDF draws each table cell as its own positioned run, so a reader that
    emits operators in order turns one row into three lines and loses the
    relationship between the cells. This fixture is the check that reading
    order is rebuilt from geometry.
    """
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4

    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Helvetica", 10)

    y = 760
    for i in ("one", "two", "three"):
        c.drawString(60, y, f"LEFT COLUMN line {i}")
        c.drawString(320, y, f"RIGHT COLUMN line {i}")
        y -= 18

    y -= 40
    for row in TABLE:
        x = 60
        for cell in row:
            c.drawString(x, y, cell)
            x += 130
        y -= 16
    c.save()
    return "reportlab, two columns and a table"


def pptx(path):
    import pptx as _pptx
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{MARKER} deck"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    box.text_frame.text = SENTENCE

    two = prs.slides.add_slide(prs.slide_layouts[5])
    two.shapes.title.text = "Second slide"
    two.notes_slide.notes_text_frame.text = "Speaker note on slide two."

    prs.save(path)
    return f"python-pptx {_pptx.__version__}"


def pdf_simple(path):
    """Standard Type1 font: no ToUnicode map, so this exercises the WinAnsi
    fallback that most business PDFs actually rely on."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4

    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 760, f"{MARKER} annual statement")
    c.setFont("Helvetica", 11)
    c.drawString(72, 730, SENTENCE)
    c.drawString(72, 710, "Second line on the first page.")
    c.showPage()
    c.setFont("Helvetica", 11)
    c.drawString(72, 760, "Page two content follows.")
    c.save()
    return "reportlab, Helvetica (no ToUnicode)"


def pdf_embedded(path):
    """Embedded TrueType subset, which carries a ToUnicode CMap. This is the
    path that produces mojibake if the CMap is parsed wrongly, so it is the
    one worth testing."""
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.pagesizes import A4

    candidates = [Path("C:/Windows/Fonts/arial.ttf"),
                  Path("C:/Windows/Fonts/calibri.ttf"),
                  Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")]
    font = next((f for f in candidates if f.exists()), None)
    if not font:
        raise RuntimeError("no TrueType font available to embed")

    pdfmetrics.registerFont(TTFont("Embedded", str(font)))
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Embedded", 14)
    c.drawString(72, 760, f"{MARKER} embedded font test")
    c.drawString(72, 735, SENTENCE)
    c.drawString(72, 710, "Curly quotes: \u201cquoted\u201d and an en dash \u2013 here.")
    c.save()
    return f"reportlab, embedded {font.name}"


def pdf_scanned(path):
    """A page with an image and no text layer: the parser must say so rather
    than returning an empty string and pretending it worked."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    import io
    import struct
    import zlib

    # a 2x2 PNG, written by hand so no imaging library is needed
    def chunk(tag, data):
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    raw = b"".join(b"\x00" + b"\xc8\x64\x32" * 2 for _ in range(2))
    png = (b"\x89PNG\r\n\x1a\n"
           + chunk(b"IHDR", struct.pack(">IIBBBBB", 2, 2, 8, 2, 0, 0, 0))
           + chunk(b"IDAT", zlib.compress(raw))
           + chunk(b"IEND", b""))

    c = canvas.Canvas(str(path), pagesize=A4)
    c.drawImage(ImageReader(io.BytesIO(png)), 72, 600, width=200, height=200)
    c.save()
    return "reportlab, image only"


def csv_file(path):
    path.write_text(
        "Site,Status,Note\n"
        f'Northam,compliant,"Reviewed, signed and {MARKER}"\n'
        'Geraldton,at risk,"Quote ""inside"" a field"\n',
        encoding="utf-8")
    return "hand-written, RFC 4180 quoting"


BUILDERS = [
    ("sample.docx", docx),
    ("sample.xlsx", xlsx),
    ("sample.pptx", pptx),
    ("simple.pdf", pdf_simple),
    ("embedded.pdf", pdf_embedded),
    ("scanned.pdf", pdf_scanned),
    ("layout.pdf", pdf_layout),
    ("sample.csv", csv_file),
]


def main():
    if OUT.exists():
        shutil.rmtree(OUT)
    OUT.mkdir(parents=True)

    made, skipped = 0, 0
    print("Generating parser fixtures")
    for name, build in BUILDERS:
        target = OUT / name
        try:
            note = build(target)
            size = target.stat().st_size
            print(f"  {name:<16} {size:>8,} bytes   {note}")
            made += 1
        except Exception as err:                     # noqa: BLE001
            print(f"  {name:<16} SKIPPED — {err}")
            skipped += 1
    print(f"Done. {made} written, {skipped} skipped.")
    return 0 if made else 1


if __name__ == "__main__":
    sys.exit(main())
